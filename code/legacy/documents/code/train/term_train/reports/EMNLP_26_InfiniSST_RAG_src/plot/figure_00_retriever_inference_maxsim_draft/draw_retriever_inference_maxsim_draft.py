#!/usr/bin/env python3
"""Draw an editable four-panel retriever inference figure draft."""

from __future__ import annotations

from pathlib import Path

import matplotlib.pyplot as plt
from matplotlib.patches import Circle, FancyArrowPatch, FancyBboxPatch, Rectangle
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


OUT_DIR = Path(__file__).resolve().parent
BASE = "retriever_inference_maxsim_draft"

INK = "111827"
RED = "9B111E"
BLUE = "12A8E8"
BLUE_LIGHT = "D9ECFA"
GREEN = "BCE8A9"
GREEN_DARK = "437A35"
ORANGE = "F6B28D"
ORANGE_DARK = "B55E29"
GRAY = "EEEEEE"
GRAY_DARK = "6B7280"
LINE = "263238"


def rgb(value: str) -> RGBColor:
    return RGBColor.from_string(value)


def add_text(slide, x, y, w, h, text, size=12, color=INK, bold=False, align=PP_ALIGN.CENTER):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.clear()
    tf.margin_left = Inches(0.02)
    tf.margin_right = Inches(0.02)
    tf.margin_top = Inches(0)
    tf.margin_bottom = Inches(0)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.name = "Aptos"
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = rgb(color)
    return tb


def add_box(slide, x, y, w, h, text="", fill="FFFFFF", line=LINE, size=11, bold=False, radius=True):
    shape_type = MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE if radius else MSO_AUTO_SHAPE_TYPE.RECTANGLE
    box = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
    box.fill.solid()
    box.fill.fore_color.rgb = rgb(fill)
    box.line.color.rgb = rgb(line)
    box.line.width = Pt(1.0)
    if text:
        tf = box.text_frame
        tf.clear()
        tf.margin_left = Inches(0.04)
        tf.margin_right = Inches(0.04)
        tf.margin_top = Inches(0.02)
        tf.margin_bottom = Inches(0.02)
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = text
        r.font.name = "Aptos"
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = rgb(INK)
    return box


def add_arrow(slide, x1, y1, x2, y2, color=LINE, width=1.4):
    line = slide.shapes.add_connector(1, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    line.line.color.rgb = rgb(color)
    line.line.width = Pt(width)
    line.line.end_arrowhead = True
    return line


def add_panel(slide, x, y, w, h, title):
    add_box(slide, x, y, w, h, "", GRAY, "FFFFFF", radius=True)
    add_text(slide, x + 0.12, y + 0.12, w - 0.24, 0.42, title, 12, INK, True)


def add_cross(slide, x, y, w, h):
    l1 = slide.shapes.add_connector(1, Inches(x), Inches(y), Inches(x + w), Inches(y + h))
    l2 = slide.shapes.add_connector(1, Inches(x + w), Inches(y), Inches(x), Inches(y + h))
    for line in (l1, l2):
        line.line.color.rgb = rgb(LINE)
        line.line.width = Pt(1.8)


def draw_pptx() -> None:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_text(slide, 0.35, 0.12, 4.1, 0.45, "Retriever Architecture", 22, RED, False, PP_ALIGN.LEFT)
    add_text(slide, 3.25, 0.78, 6.9, 0.32, "Timeline-Aware MaxSim Term Retrieval", 16, INK, True)

    px = [0.35, 3.55, 6.75, 9.95]
    py, pw, ph = 1.22, 2.95, 5.55
    titles = [
        "1. Timeline-Aware Audio\nContext",
        "2. Speech Window\nConstruction",
        "3. MaxSim Retrieval against\nGlossary Bank",
        "4. Filtering and\nTerm-Map Output",
    ]
    for x, title in zip(px, titles):
        add_panel(slide, x, py, pw, ph, title)

    # Panel 1: timeline.
    x = px[0]
    add_text(slide, x + 0.15, py + 1.10, 1.05, 0.22, "speech timeline", 10, INK, False, PP_ALIGN.LEFT)
    y = py + 2.05
    add_arrow(slide, x + 0.18, y + 0.25, x + 2.68, y + 0.25, LINE, 1.1)
    add_text(slide, x + 0.18, y - 0.32, 0.65, 0.25, "previous\nchunk", 9.5, INK)
    add_text(slide, x + 2.18, y - 0.32, 0.50, 0.25, "next\nchunk", 9.5, INK)
    add_box(slide, x + 0.88, y - 0.38, 0.88, 0.95, "", "D8D8D8", GRAY_DARK, radius=False)
    add_box(slide, x + 1.76, y - 0.95, 0.62, 1.52, "cᵢ\ncurrent\nchunk", BLUE, "067BAE", 9, True, radius=False)
    b = slide.shapes.add_connector(1, Inches(x + 1.76), Inches(y - 1.18), Inches(x + 1.76), Inches(y + 0.86))
    b.line.color.rgb = rgb(LINE)
    b.line.width = Pt(1.0)
    b.line.dash_style = 4
    add_text(slide, x + 1.78, y - 1.36, 0.75, 0.25, "current-step\nboundary", 8, INK)
    add_text(slide, x + 0.40, y + 0.45, 0.45, 0.25, "cᵢ₋₁", 9, INK)
    add_text(slide, x + 2.43, y + 0.45, 0.42, 0.25, "cᵢ₊₁", 9, INK)
    add_text(slide, x + 0.95, y + 0.13, 0.75, 0.25, "1.92s\nlook-back", 9, INK)
    add_text(slide, x + 0.55, py + 4.18, 1.95, 0.45, "audio context\n aᵢ = look-back + current chunk", 10, INK)
    add_arrow(slide, px[0] + 2.66, py + 3.08, px[1] + 0.30, py + 3.08, LINE)

    # Panel 2: window construction.
    x = px[1]
    add_box(slide, x + 0.55, py + 0.78, 1.35, 0.52, "Speech\nencoder", GREEN, GREEN_DARK, 12, False)
    add_arrow(slide, x + 1.22, py + 1.30, x + 1.22, py + 1.78, LINE)
    for i, label in enumerate(["h₁", "h₂", "...", "hₘ"]):
        cx = x + 0.55 + i * 0.48
        circ = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(cx), Inches(py + 1.74), Inches(0.34), Inches(0.34))
        circ.fill.solid()
        circ.fill.fore_color.rgb = rgb("A9D5FF" if label != "..." else "FFFFFF")
        circ.line.color.rgb = rgb("4B83B7" if label != "..." else GRAY)
        if label != "...":
            add_text(slide, cx + 0.02, py + 1.79, 0.30, 0.20, label, 8.5, INK)
        else:
            add_text(slide, cx, py + 1.80, 0.34, 0.20, "...", 13, INK, True)
    for dx, dy, ww in [(0.35, 2.62, 0.95), (0.72, 2.82, 1.05), (0.95, 3.05, 1.25), (1.30, 3.32, 1.12)]:
        add_box(slide, x + dx, py + dy, ww, 0.30, "", BLUE_LIGHT, "6C95C8", radius=False)
    add_arrow(slide, x + 0.72, py + 3.55, x + 0.68, py + 3.14, GRAY_DARK, 1.0)
    add_arrow(slide, x + 1.05, py + 3.65, x + 1.22, py + 3.22, GRAY_DARK, 1.0)
    add_text(slide, x + 0.25, py + 3.58, 0.85, 0.42, "may start,\ni.e. overlap cᵢ", 8.5, INK, False, PP_ALIGN.LEFT)
    add_text(slide, x + 1.00, py + 4.28, 1.55, 0.38, "multi-scale\nspeech windows", 10.5, INK)
    add_arrow(slide, x + 2.72, py + 1.00, x + 3.08, py + 1.00, LINE)
    add_arrow(slide, x + 2.72, py + 2.00, x + 3.08, py + 2.00, LINE)
    add_arrow(slide, x + 2.72, py + 3.18, x + 3.08, py + 3.18, LINE)

    # Panel 3: heatmap and glossary.
    x = px[2]
    add_text(slide, x + 0.20, py + 1.00, 1.05, 0.35, "Cosine Similarity", 10.5, INK)
    hx, hy = x + 0.28, py + 1.55
    colors = [
        ["64B87A", "A8CE81", "F3DD7C", "F4A65D"],
        ["76C389", "B8D38C", "F2D36C", "F28C50"],
        ["78C58A", "ABD184", "EAD86E", "F0A25F"],
        ["75C187", "A4CD7D", "E5D56E", "F4B56D"],
        ["6CBA7D", "9FC978", "E0CF6D", "F2A960"],
    ]
    for r in range(5):
        for c in range(4):
            add_box(slide, hx + c * 0.33, hy + r * 0.33, 0.33, 0.33, "", colors[r][c], colors[r][c], radius=False)
    add_box(slide, hx + 0.63, hy + 0.96, 0.78, 0.36, "max(...)", GREEN, GREEN_DARK, 10, False)
    add_text(slide, x + 0.22, py + 4.40, 2.15, 0.52, "MaxSim:\nscore(eⱼ) = max over\nspeech windows zᵀ gⱼ", 10, INK)
    gx = x + 2.08
    add_text(slide, gx - 0.20, py + 0.85, 0.9, 0.45, "Pre-encoded\nglossary bank", 9, INK)
    for i, label in enumerate(["g₁", "g₂", "...", "gₙ"]):
        if label == "...":
            add_text(slide, gx + 0.10, py + 2.70, 0.35, 0.35, "...", 16, INK, True)
            continue
        yy = py + 1.48 + i * 0.48 if i < 2 else py + 3.65
        add_box(slide, gx, yy, 0.48, 0.42, label, ORANGE, ORANGE_DARK, 10, False, radius=False)
    for yy in [py + 1.72, py + 2.20, py + 4.00]:
        add_arrow(slide, hx + 1.32, hy + 0.98, gx, yy, ORANGE_DARK, 0.9)
        add_arrow(slide, gx, yy, hx + 1.30, hy + 1.20, GREEN_DARK, 0.9)
    add_arrow(slide, x + 2.72, py + 2.25, x + 3.05, py + 2.25, LINE)

    # Panel 4: filtering and output.
    x = px[3]
    add_box(slide, x + 0.28, py + 0.72, 2.35, 0.82, "logical:\n- overlap current chunk\n- score ≥ τ\n- top-K", GREEN, GREEN_DARK, 11, False)
    add_arrow(slide, x + 1.45, py + 1.55, x + 1.45, py + 2.15, LINE)
    add_box(slide, x + 0.62, py + 2.20, 0.82, 0.48, "", BLUE_LIGHT, "6C95C8", radius=False)
    add_box(slide, x + 1.98, py + 2.22, 0.46, 0.46, "", BLUE_LIGHT, "6C95C8", radius=False)
    add_cross(slide, x + 0.58, py + 2.14, 0.90, 0.60)
    add_cross(slide, x + 1.94, py + 2.16, 0.54, 0.55)
    add_arrow(slide, x + 1.45, py + 2.78, x + 1.45, py + 3.10, LINE)
    add_text(slide, x + 0.32, py + 3.15, 2.25, 0.28, "step-local term map Gᵢ", 12, INK)
    add_box(slide, x + 0.28, py + 3.47, 2.35, 0.92, "source term -> target term\nsource term -> target term\n...", ORANGE, ORANGE_DARK, 10, False)
    add_text(slide, x + 0.12, py + 4.82, 2.70, 0.45, "stateless per step: Gᵢ is recomputed\nfrom aᵢ and the glossary bank", 8.6, INK)

    prs.save(OUT_DIR / f"{BASE}.pptx")


def box(ax, x, y, w, h, text="", fc="FFFFFF", ec=LINE, size=10, weight="normal", radius=0.04):
    patch = FancyBboxPatch(
        (x, y),
        w,
        h,
        boxstyle=f"round,pad=0.01,rounding_size={radius}",
        facecolor=f"#{fc}",
        edgecolor=f"#{ec}",
        linewidth=1.0,
    )
    ax.add_patch(patch)
    if text:
        ax.text(x + w / 2, y + h / 2, text, ha="center", va="center", fontsize=size, weight=weight, color=f"#{INK}")
    return patch


def arrow(ax, start, end, color=LINE, lw=1.1):
    ax.add_patch(FancyArrowPatch(start, end, arrowstyle="-|>", mutation_scale=10, linewidth=lw, color=f"#{color}"))


def cross(ax, x, y, w, h):
    ax.plot([x, x + w], [y, y + h], color=f"#{LINE}", lw=1.4)
    ax.plot([x + w, x], [y, y + h], color=f"#{LINE}", lw=1.4)


def draw_preview() -> None:
    fig, ax = plt.subplots(figsize=(13.33, 7.5))
    ax.set_xlim(0, 13.33)
    ax.set_ylim(0, 7.5)
    ax.axis("off")

    ax.text(0.35, 7.15, "Retriever Architecture", ha="left", va="center", fontsize=22, color=f"#{RED}")
    ax.text(6.67, 6.60, "Timeline-Aware MaxSim Term Retrieval", ha="center", va="center", fontsize=16, weight="bold", color=f"#{INK}")

    px = [0.35, 3.55, 6.75, 9.95]
    py, pw, ph = 0.92, 2.95, 5.55
    titles = [
        "1. Timeline-Aware Audio\nContext",
        "2. Speech Window\nConstruction",
        "3. MaxSim Retrieval against\nGlossary Bank",
        "4. Filtering and\nTerm-Map Output",
    ]
    for x, title in zip(px, titles):
        box(ax, x, py, pw, ph, "", GRAY, "FFFFFF", radius=0.05)
        ax.text(x + pw / 2, py + ph - 0.42, title, ha="center", va="center", fontsize=12, weight="bold", color=f"#{INK}")

    # Panel 1.
    x = px[0]
    y = py + 2.85
    ax.text(x + 0.15, py + 4.05, "speech timeline", ha="left", fontsize=10, color=f"#{INK}")
    arrow(ax, (x + 0.18, y + 0.25), (x + 2.72, y + 0.25))
    ax.text(x + 0.36, y + 0.67, "previous\nchunk", ha="center", fontsize=9.5)
    ax.text(x + 2.40, y + 0.67, "next\nchunk", ha="center", fontsize=9.5)
    ax.add_patch(Rectangle((x + 0.88, y - 0.38), 0.88, 0.95, facecolor="#D8D8D8", edgecolor=f"#{GRAY_DARK}"))
    ax.add_patch(Rectangle((x + 1.76, y - 0.95), 0.62, 1.52, facecolor=f"#{BLUE}", edgecolor="#067BAE"))
    ax.plot([x + 1.76, x + 1.76], [y - 1.20, y + 0.90], "--", color=f"#{LINE}", lw=1.0)
    ax.text(x + 2.02, y + 1.02, "current-step\nboundary", ha="center", fontsize=8)
    ax.text(x + 0.55, y - 0.05, "$c_{i-1}$", ha="center", fontsize=9)
    ax.text(x + 2.65, y - 0.05, "$c_{i+1}$", ha="center", fontsize=9)
    ax.text(x + 1.33, y - 0.02, "1.92s\nlook-back", ha="center", fontsize=9)
    ax.text(x + 2.07, y - 0.15, "$c_i$\ncurrent\nchunk", ha="center", va="center", fontsize=9, color="white", weight="bold")
    ax.text(x + 1.48, py + 0.72, "audio context\n$a_i$ = look-back + current chunk", ha="center", fontsize=10)
    arrow(ax, (px[0] + 2.66, py + 3.08), (px[1] + 0.30, py + 3.08))

    # Panel 2.
    x = px[1]
    box(ax, x + 0.55, py + 4.05, 1.35, 0.52, "Speech\nencoder", GREEN, GREEN_DARK, 12)
    arrow(ax, (x + 1.22, py + 4.05), (x + 1.22, py + 3.58))
    for i, label in enumerate(["$h_1$", "$h_2$", "...", "$h_M$"]):
        cx, cy = x + 0.55 + i * 0.48, py + 3.25
        if label == "...":
            ax.text(cx + 0.17, cy + 0.17, "...", ha="center", va="center", fontsize=14, weight="bold")
        else:
            circ = Circle((cx + 0.17, cy + 0.17), 0.17, facecolor="#A9D5FF", edgecolor="#4B83B7")
            ax.add_patch(circ)
            ax.text(cx + 0.17, cy + 0.17, label, ha="center", va="center", fontsize=8.5, style="italic")
    for dx, dy, ww in [(0.35, 2.42, 0.95), (0.72, 2.22, 1.05), (0.95, 1.99, 1.25), (1.30, 1.72, 1.12)]:
        ax.add_patch(Rectangle((x + dx, py + dy), ww, 0.30, facecolor=f"#{BLUE_LIGHT}", edgecolor="#6C95C8"))
    arrow(ax, (x + 0.72, py + 1.55), (x + 0.70, py + 2.35), GRAY_DARK, 0.9)
    arrow(ax, (x + 1.06, py + 1.44), (x + 1.20, py + 2.10), GRAY_DARK, 0.9)
    ax.text(x + 0.30, py + 1.28, "may start, i.e.\noverlap $c_i$", ha="left", fontsize=8.5)
    ax.text(x + 1.62, py + 0.66, "multi-scale\nspeech windows", ha="center", fontsize=10.5)
    for yy in [py + 4.30, py + 3.40, py + 2.15]:
        arrow(ax, (x + 2.73, yy), (x + 3.08, yy))

    # Panel 3.
    x = px[2]
    ax.text(x + 0.75, py + 4.25, "Cosine Similarity", ha="center", fontsize=10.5)
    hx, hy = x + 0.28, py + 2.30
    colors = [
        ["64B87A", "A8CE81", "F3DD7C", "F4A65D"],
        ["76C389", "B8D38C", "F2D36C", "F28C50"],
        ["78C58A", "ABD184", "EAD86E", "F0A25F"],
        ["75C187", "A4CD7D", "E5D56E", "F4B56D"],
        ["6CBA7D", "9FC978", "E0CF6D", "F2A960"],
    ]
    for r in range(5):
        for c in range(4):
            ax.add_patch(Rectangle((hx + c * 0.33, hy + (4 - r) * 0.33), 0.33, 0.33, facecolor=f"#{colors[r][c]}", edgecolor=f"#{colors[r][c]}"))
    box(ax, hx + 0.63, hy + 0.66, 0.78, 0.36, "max(...)", GREEN, GREEN_DARK, 10)
    ax.text(x + 1.15, py + 0.60, "MaxSim:\n$score(e_j)$ = max over\nspeech windows $z^T g_j$", ha="center", fontsize=10)
    gx = x + 2.08
    ax.text(gx + 0.20, py + 4.20, "Pre-encoded\nglossary bank", ha="center", fontsize=9)
    for yy, label in [(py + 3.47, "$g_1$"), (py + 2.99, "$g_2$"), (py + 1.35, "$g_N$")]:
        ax.add_patch(Rectangle((gx, yy), 0.48, 0.42, facecolor=f"#{ORANGE}", edgecolor=f"#{ORANGE_DARK}"))
        ax.text(gx + 0.24, yy + 0.21, label, ha="center", va="center", fontsize=10)
    ax.text(gx + 0.24, py + 2.22, "...", ha="center", fontsize=18, weight="bold")
    for yy in [py + 3.68, py + 3.20, py + 1.56]:
        arrow(ax, (hx + 1.32, hy + 1.0), (gx, yy), ORANGE_DARK, 0.8)
        arrow(ax, (gx, yy), (hx + 1.30, hy + 1.2), GREEN_DARK, 0.8)
    arrow(ax, (x + 2.73, py + 3.15), (x + 3.05, py + 3.15))

    # Panel 4.
    x = px[3]
    box(ax, x + 0.28, py + 4.08, 2.35, 0.82, "logical:\n- overlap current chunk\n- score >= $\\tau$\n- top-K", GREEN, GREEN_DARK, 11)
    arrow(ax, (x + 1.45, py + 4.05), (x + 1.45, py + 3.45))
    ax.add_patch(Rectangle((x + 0.62, py + 2.95), 0.82, 0.48, facecolor=f"#{BLUE_LIGHT}", edgecolor="#6C95C8"))
    ax.add_patch(Rectangle((x + 1.98, py + 2.97), 0.46, 0.46, facecolor=f"#{BLUE_LIGHT}", edgecolor="#6C95C8"))
    cross(ax, x + 0.58, py + 2.90, 0.90, 0.60)
    cross(ax, x + 1.94, py + 2.92, 0.54, 0.55)
    arrow(ax, (x + 1.45, py + 2.86), (x + 1.45, py + 2.55))
    ax.text(x + 1.45, py + 2.34, "step-local term map $G_i$", ha="center", fontsize=12)
    box(ax, x + 0.28, py + 1.30, 2.35, 0.92, "source term -> target term\nsource term -> target term\n...", ORANGE, ORANGE_DARK, 10)
    ax.text(x + 1.45, py + 0.28, "stateless per step: $G_i$ is recomputed\nfrom $a_i$ and the glossary bank", ha="center", fontsize=8.6, style="italic")

    fig.savefig(OUT_DIR / f"{BASE}.png", dpi=300, bbox_inches="tight")
    fig.savefig(OUT_DIR / f"{BASE}.pdf", bbox_inches="tight")
    fig.savefig(OUT_DIR / f"{BASE}.svg", bbox_inches="tight")
    plt.close(fig)


def main() -> int:
    draw_pptx()
    draw_preview()
    print(f"Wrote {OUT_DIR / (BASE + '.pptx')}")
    print(f"Wrote {OUT_DIR / (BASE + '.pdf')}")
    print(f"Wrote {OUT_DIR / (BASE + '.png')}")
    print(f"Wrote {OUT_DIR / (BASE + '.svg')}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
